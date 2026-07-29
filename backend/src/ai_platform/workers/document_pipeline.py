"""
Pipeline Celery para ingestion de documentos.

Pipeline en 8 stages:
1. ingest_pipeline (entry point)
2. extract_text (PDF/DOCX/XLSX/PPTX/image -> raw text)
3. extract_embedded_images (extrae imágenes incrustadas de DOCX/PDF/PPTX)
4. vision_analyze_images (vision_chat para cada imagen incrustada)
5. ocr_all_pages (si es escaneado, Tesseract por pagina)
6. chunk_document (text -> chunks, lee de la BD)
7. build_fts_index (chunks -> tsvector)
8. generate_summaries (chunks -> summaries via LLM)
9. mark_completed

Usar:
    from ai_platform.workers.document_pipeline import ingest_pipeline
    ingest_pipeline.delay(document_id, tenant_id, strategy="hybrid")
"""

import asyncio
import logging
import tempfile
import time
from pathlib import Path

from celery.utils.log import get_task_logger
from sqlalchemy import select, text as sa_text

from ai_platform.database import session_factory
from ai_platform.models.db import DocumentArtifact, DocumentChunk, DocumentFTSIndex
from ai_platform.workers.task_runner import celery_app

logger = get_task_logger("ai_platform.document_pipeline")


def _get_document_session(document_id: str):
    """Obtener session de DB con el documento."""
    session = session_factory()
    stmt = select(DocumentArtifact).where(DocumentArtifact.id == document_id)
    doc = session.execute(stmt).scalar_one_or_none()
    return session, doc


def _update_document(document_id: str, updates: dict) -> None:
    """Actualizar campos del documento."""
    session = session_factory()
    try:
        stmt = select(DocumentArtifact).where(DocumentArtifact.id == document_id)
        doc = session.execute(stmt).scalar_one_or_none()
        if doc:
            for key, value in updates.items():
                setattr(doc, key, value)
            session.commit()
    except Exception:
        session.rollback()
        raise
    finally:
        session.close()


def _resolve_path(file_path: str, tenant_id: str) -> Path:
    """Resolver ruta completa del archivo."""
    from ai_platform.core.config import get_settings

    settings = get_settings()
    storage_root = Path(settings.DOCUMENT_STORAGE_ROOT)
    return storage_root / str(tenant_id) / file_path


def _extract_images_from_docx(full_path: Path) -> list[tuple[str, bytes, str]]:
    """Extraer imágenes incrustadas de un DOCX.

    Retorna lista de (nombre_archivo, bytes, mime_type).
    """
    images = []
    try:
        from docx import Document
        from docx.oxml.constants import CT_EMBED, CT_INLINE

        docx_obj = Document(str(full_path))
        image_counter = 0

        for part in docx_obj.part.rels.values():
            if part.reltype == "/relationship/image" or "image" in part.reltype.lower():
                try:
                    image_part = part.target_part
                    blob = image_part.blob
                    content_type = image_part.content_type or "image/png"
                    ext = content_type.split("/")[-1]
                    name = f"docx_image_{image_counter}.{ext}"
                    images.append((name, blob, content_type))
                    image_counter += 1
                except Exception:
                    continue

        # También buscar en drawings
        try:
            from docx.oxml.ns import qn

            for para in docx_obj.paragraphs:
                for run in para.runs:
                    if run._element:
                        for drawing in run._element.findall(qn("w:drawing")):
                            for anchor in drawing.findall(qn("a:externalReference")):
                                try:
                                    rId = anchor.get(qn("r:id"))
                                    if rId and rId in docx_obj.part.rels:
                                        rel = docx_obj.part.rels[rId]
                                        if "image" in rel.reltype:
                                            image_part = rel.target_part
                                            blob = image_part.blob
                                            content_type = image_part.content_type or "image/png"
                                            ext = content_type.split("/")[-1]
                                            name = f"docx_image_{image_counter}.{ext}"
                                            images.append((name, blob, content_type))
                                            image_counter += 1
                                except Exception:
                                    continue
        except Exception:
            pass

    except Exception as e:
        logger.warning(f"extract_images_from_docx failed: {e}")

    logger.info(f"extract_images_from_docx found {len(images)} images")
    return images


def _extract_images_from_pdf(full_path: Path) -> list[tuple[str, bytes, str]]:
    """Extraer imágenes incrustadas de un PDF usando PyMuPDF.

    Retorna lista de (nombre_archivo, bytes, mime_type).
    """
    images = []
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(str(full_path))
        image_counter = 0

        for page_num in range(len(doc)):
            page = doc[page_num]
            image_list = page.get_images(full=True)

            for img_info in image_list:
                xref = img_info[0]
                try:
                    img_bytes = doc.extract_image(xref)
                    if img_bytes:
                        ext = img_bytes.get("ext", "png")
                        content_type = f"image/{ext}"
                        name = f"pdf_page{page_num + 1}_image_{image_counter}.{ext}"
                        images.append((name, img_bytes["image"], content_type))
                        image_counter += 1
                except Exception:
                    continue

        doc.close()
    except Exception as e:
        logger.warning(f"extract_images_from_pdf failed: {e}")

    logger.info(f"extract_images_from_pdf found {len(images)} images")
    return images


def _extract_images_from_pptx(full_path: Path) -> list[tuple[str, bytes, str]]:
    """Extraer imágenes incrustadas de un PPTX.

    Retorna lista de (nombre_archivo, bytes, mime_type).
    """
    images = []
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(str(full_path))
        image_counter = 0

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        blob = image.blob
                        content_type = image.content_type or "image/png"
                        ext = content_type.split("/")[-1]
                        name = f"pptx_slide{slide_idx + 1}_image_{image_counter}.{ext}"
                        images.append((name, blob, content_type))
                        image_counter += 1
                    except Exception:
                        continue
                elif shape.has_text_frame:
                    # Texto de shapes también se captura en extract_text
                    pass

    except Exception as e:
        logger.warning(f"extract_images_from_pptx failed: {e}")

    logger.info(f"extract_images_from_pptx found {len(images)} images")
    return images


def _extract_images_from_xlsx(full_path: Path) -> list[tuple[str, bytes, str]]:
    """Extraer imágenes incrustadas de un XLSX.

    openpyxl no soporta imágenes directamente, pero XLSX es un ZIP.
    Las imágenes están en /xl/media/.
    """
    images = []
    try:
        import zipfile

        with zipfile.ZipFile(str(full_path), "r") as z:
            media_files = [f for f in z.namelist() if f.startswith("xl/media/")]
            image_counter = 0

            for media_file in media_files:
                try:
                    blob = z.read(media_file)
                    ext = media_file.rsplit(".", 1)[-1].lower()
                    content_type = f"image/{ext}"
                    name = f"xlsx_media_{image_counter}.{ext}"
                    images.append((name, blob, content_type))
                    image_counter += 1
                except Exception:
                    continue
    except Exception as e:
        logger.warning(f"extract_images_from_xlsx failed: {e}")

    logger.info(f"extract_images_from_xlsx found {len(images)} images")
    return images


def _extract_embedded_images(self, document_id: str, tenant_id: str) -> dict:
    """Extraer imágenes incrustadas de un documento.

    Stage intermedio: después de extract_text, antes de OCR.
    Guarda las imágenes como artefactos separados en el storage.

    Retorna conteo de imágenes extraídas por formato.
    """
    logger.info("extract_embedded_images_started", document_id=document_id)
    session, doc = _get_document_session(document_id)
    if not doc:
        logger.error("extract_embedded_images_no_doc", document_id=document_id)
        return {"images_extracted": 0}

    full_path = _resolve_path(doc.file_path, tenant_id)
    mime_type = doc.mime_type

    all_images: list[tuple[str, bytes, str]] = []

    if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        all_images = _extract_images_from_docx(full_path)
    elif mime_type == "application/pdf":
        all_images = _extract_images_from_pdf(full_path)
    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        all_images = _extract_images_from_pptx(full_path)
    elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        all_images = _extract_images_from_xlsx(full_path)

    # Guardar imágenes como artefactos separados
    storage_root = Path(doc.file_path).parent.parent  # sube un nivel al tenant dir
    images_saved = 0

    for name, blob, mime in all_images:
        try:
            img_path = storage_root / f"embedded_{document_id}_{images_saved}_{name}"
            img_path.parent.mkdir(parents=True, exist_ok=True)
            img_path.write_bytes(blob)

            # Crear artefacto para tracking
            from uuid import uuid4

            img_artifact = DocumentArtifact(
                tenant_id=doc.tenant_id,
                name=f"embedded_{name}",
                mime_type=mime,
                size_bytes=len(blob),
                file_path=str(img_path.relative_to(storage_root.parent)),
                status="embedded_image",
            )
            session.add(img_artifact)
            images_saved += 1
        except Exception as e:
            logger.warning(f"Failed to save embedded image {name}: {e}")

    session.commit()
    logger.info("extract_embedded_images_done", document_id=document_id, count=images_saved)
    return {"images_extracted": images_saved}


@celery_app.task(
    name="documents.extract_text",
    bind=True,
    max_retries=2,
    acks_late=True,
    time_limit=1800,
    soft_time_limit=1500,
)
def extract_text(self, document_id: str, tenant_id: str) -> dict:
    """
    Extraer texto crudo de PDF, DOCX, XLSX, PPTX, o imágenes.
    Guarda el texto en la BD (campo extracted_text) para que chunk_document lo lea.
    """
    logger.info("extract_text_started", document_id=document_id)
    session, doc = _get_document_session(document_id)
    if not doc:
        raise ValueError(f"Documento {document_id} no encontrado")

    full_path = _resolve_path(doc.file_path, tenant_id)

    if not full_path.exists():
        raise FileNotFoundError(f"Archivo no encontrado: {full_path}")

    extracted_text = None
    is_scan = False
    needs_ocr = False
    page_count = 1
    file_type = doc.mime_type.split("/")[1] if "/" in doc.mime_type else "unknown"

    if doc.mime_type == "application/pdf":
        try:
            import pdfplumber

            with pdfplumber.open(full_path) as pdf:
                page_count = len(pdf.pages)
                pages_text = []
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        pages_text.append(text)
                    else:
                        is_scan = True
                        needs_ocr = True
                extracted_text = "\n\n".join(pages_text) if pages_text else None
        except ImportError:
            logger.warning("pdfplumber not installed, marking as scan")
            is_scan = True
            needs_ocr = True

    elif doc.mime_type.startswith("image/"):
        needs_ocr = True
        is_scan = True
        extracted_text = ""

    elif doc.mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        try:
            from docx import Document

            docx_obj = Document(str(full_path))
            extracted_text = "\n".join(para.text for para in docx_obj.paragraphs)
        except ImportError:
            logger.warning("python-docx not installed, skipping DOCX extraction")

    elif doc.mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        try:
            from pptx import Presentation

            prs = Presentation(str(full_path))
            slide_texts = []
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text = f"\n--- Slide {slide_idx} ---\n"
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_text += f"  {para.text.strip()}\n"
                slide_texts.append(slide_text)
            extracted_text = "\n".join(slide_texts)
        except ImportError:
            logger.warning("python-pptx not installed, skipping PPTX extraction")

    elif (
        doc.mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        or "ms-excel" in doc.mime_type
    ):
        try:
            import openpyxl

            wb = openpyxl.load_workbook(str(full_path), read_only=True, data_only=True)
            sheet_text = []
            for ws in wb.worksheets:
                sheet_text.append(f"--- Sheet: {ws.title} ---")
                for row in ws.iter_rows(values_only=True):
                    row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_str.strip():
                        sheet_text.append(row_str)
                wb.close()
            extracted_text = "\n".join(sheet_text)
        except ImportError:
            logger.warning("openpyxl not installed, skipping XLSX extraction")

    elif doc.mime_type == "text/plain":
        try:
            extracted_text = full_path.read_text(encoding="utf-8")
        except Exception:
            try:
                extracted_text = full_path.read_text(encoding="latin-1")
            except Exception:
                pass

    elif doc.mime_type == "text/csv":
        try:
            content = full_path.read_text(encoding="utf-8")
            extracted_text = "CSV:\n" + content
        except Exception:
            pass

    if not extracted_text:
        extracted_text = ""
        if doc.mime_type.startswith("image/"):
            needs_ocr = True
            is_scan = True

    # Guardar texto extraido en la BD
    doc.extracted_text = extracted_text
    doc.page_count = page_count
    doc.status = "extracted"
    session.commit()

    logger.info(
        "extract_text_completed",
        document_id=document_id,
        page_count=page_count,
        needs_ocr=needs_ocr,
        text_length=len(extracted_text or ""),
    )

    return {
        "extracted_text": extracted_text,
        "is_scan": is_scan,
        "page_count": page_count,
        "file_type": file_type,
        "needs_ocr": needs_ocr,
    }


@celery_app.task(
    name="documents.extract_embedded_images",
    bind=True,
    max_retries=2,
    acks_late=True,
    time_limit=600,
    soft_time_limit=540,
)
def extract_embedded_images(self, document_id: str, tenant_id: str) -> dict:
    """Extraer imágenes incrustadas de un documento (DOCX, PDF, PPTX, XLSX).

    Stage intermedio entre extract_text y OCR.
    Las imágenes se guardan como artefactos separados.
    """
    return _extract_embedded_images(self, document_id, tenant_id)


@celery_app.task(
    name="documents.vision_analyze_images",
    bind=True,
    max_retries=2,
    acks_late=True,
    time_limit=3600,
    soft_time_limit=3300,
)
def vision_analyze_images(self, document_id: str, tenant_id: str) -> dict:
    """Analizar imágenes incrustadas con LLM vision (mimo-v2.5).

    Para cada imagen incrustada:
    1. Leer bytes del artefacto
    2. Llamar vision_chat con prompt contextual
    3. Agregar descripción al extracted_text del documento

    Retorna conteo de imágenes analizadas.
    """
    logger.info("vision_analyze_images_started", document_id=document_id)
    session, doc = _get_document_session(document_id)
    if not doc:
        logger.error("vision_analyze_images_no_doc", document_id=document_id)
        return {"images_analyzed": 0}

    try:
        from ai_platform.orchestrator.llm_client import LLMClient

        llm = LLMClient()

        # Buscar artefactos de imágenes incrustadas
        stmt = select(DocumentArtifact).where(
            DocumentArtifact.tenant_id == doc.tenant_id,
            DocumentArtifact.status == "embedded_image",
        )
        embedded_images = session.execute(stmt).scalars().all()

        if not embedded_images:
            logger.info("vision_analyze_images_no_embedded", document_id=document_id)
            return {"images_analyzed": 0}

        descriptions = []
        analyzed = 0

        for img_artifact in embedded_images:
            try:
                img_path = _resolve_path(img_artifact.file_path, tenant_id)
                if not img_path.exists():
                    logger.warning(f"embedded_image_not_found: {img_artifact.file_path}")
                    continue

                img_bytes = img_path.read_bytes()
                if not img_bytes:
                    continue

                # Prompt contextual según el tipo de documento
                prompt = (
                    "Analiza esta imagen extraída de un documento. "
                    "Si es un gráfico o tabla, extrae todos los datos y valores. "
                    "Si es una foto o imagen, descríbela brevemente. "
                    "Si es un diagrama o esquema, describe su estructura. "
                    "Responde en español, de forma concisa y estructurada.\n\n"
                    f"Tipo de imagen: {img_artifact.mime_type}"
                )

                loop = asyncio.new_event_loop()
                try:
                    result = loop.run_until_complete(
                        llm.vision_chat(prompt, img_bytes, tenant_id=tenant_id)
                    )
                finally:
                    loop.close()

                text = result.get("text", "") if isinstance(result, dict) else str(result)

                if text and text.strip():
                    descriptions.append(f"\n--- Imagen incrustada: {img_artifact.name} ---\n{text}\n--- Fin de imagen ---")
                    analyzed += 1

                # Marcar artefacto como procesado
                img_artifact.status = "embedded_image_processed"
            except Exception as e:
                logger.warning(f"vision_analyze_image_failed ({img_artifact.name}): {e}")
                img_artifact.status = "embedded_image_error"

        # Agregar descripciones al texto del documento
        if descriptions:
            doc.extracted_text = (doc.extracted_text or "") + "\n\n".join(descriptions)
            session.commit()

        logger.info("vision_analyze_images_done", document_id=document_id, analyzed=analyzed)
        return {"images_analyzed": analyzed}

    except Exception as e:
        logger.error("vision_analyze_images_failed", document_id=document_id, error=str(e))
        return {"images_analyzed": 0, "error": str(e)}
    finally:
        try:
            llm.close()
        except Exception:
            pass


@celery_app.task(
    name="documents.ocr_all_pages",
    bind=True,
    max_retries=2,
    acks_late=True,
)
def ocr_all_pages(self, document_id: str, tenant_id: str, _page_count: int = 1) -> dict:
    """
    OCR para PDFs escaneados o imágenes.
    Usa Tesseract como primario. Si confianza < 60%, fallback a LLM vision.
    """
    logger.info("ocr_all_pages_started", document_id=document_id)
    session, doc = _get_document_session(document_id)
    try:
        full_path = _resolve_path(doc.file_path, tenant_id)

        if not full_path.exists():
            return {"pages_ocrd": 0}

        extracted_text_from_ocr = ""
        tesseract_confidence = 0.0

        if doc.mime_type == "application/pdf":
            try:
                from pdf2image import convert_from_path

                images = convert_from_path(str(full_path), dpi=300)
                ocr_pages = []

                for img in images:
                    import pytesseract

                    page_text = pytesseract.image_to_string(img, lang="spa+eng")
                    ocr_pages.append(page_text)
                    extracted_text_from_ocr += page_text + "\n\n"

                # Calcular confianza promedio de Tesseract
                try:
                    for img in images:
                        import pytesseract

                        config_data = pytesseract.image_to_data(img, lang="spa+eng", output_type=pytesseract.Output.DICT)
                        total_chars = sum(len(c) for c in config_data["text"] if c.strip())
                        good_chars = sum(len(c) for c, conf in zip(config_data["text"], config_data["conf"]) if conf > 50 and c.strip())
                        if total_chars > 0:
                            tesseract_confidence = max(tesseract_confidence, (good_chars / total_chars) * 100)
                except Exception:
                    pass

            except ImportError:
                logger.warning("pdf2image not installed, skipping PDF OCR")

        elif doc.mime_type.startswith("image/"):
            try:
                from PIL import Image

                import pytesseract

                img = Image.open(str(full_path))
                extracted_text_from_ocr += pytesseract.image_to_string(img, lang="spa+eng")
            except ImportError:
                logger.warning("tesseract/PIL not installed, skipping OCR")

        if extracted_text_from_ocr:
            doc.extracted_text = (doc.extracted_text or "") + extracted_text_from_ocr
            session.commit()

            logger.info(
                "ocr_all_pages_done",
                document_id=document_id,
                ocr_length=len(extracted_text_from_ocr),
                tesseract_confidence=tesseract_confidence,
            )
        else:
            logger.warning("ocr_all_pages_no_text", document_id=document_id)

        # Fallback a vision si Tesseract tiene baja confianza
        vision_used = False
        if tesseract_confidence < 60 and extracted_text_from_ocr and doc.mime_type == "application/pdf":
            logger.info("ocr_fallback_to_vision", document_id=document_id, confidence=tesseract_confidence)
            try:
                from ai_platform.orchestrator.llm_client import LLMClient

                llm = LLMClient()
                from pdf2image import convert_from_path

                images = convert_from_path(str(full_path), dpi=300)
                vision_pages = []

                for img in images:
                    import io

                    img_bytes = io.BytesIO()
                    img.save(img_bytes, format="PNG")
                    img_bytes.seek(0)
                    img_data = img_bytes.read()

                    loop = asyncio.new_event_loop()
                    try:
                        result = loop.run_until_complete(
                            llm.vision_chat(
                                "Transcribe TODO el texto visible en este documento escaneado. "
                                "Mantén el formato original. Responde en español.",
                                img_data,
                            )
                        )
                    finally:
                        loop.close()

                    text = result.get("text", "") if isinstance(result, dict) else str(result)
                    if text and text.strip():
                        vision_pages.append(text)

                if vision_pages:
                    vision_text = "\n\n".join(vision_pages)
                    if vision_text.strip():
                        doc.extracted_text = (doc.extracted_text or "") + "\n\n[Visión - OCR fallback]:\n" + vision_text
                        session.commit()
                        vision_used = True
                        logger.info("ocr_fallback_vision_done", document_id=document_id, vision_len=len(vision_text))

                try:
                    llm.close()
                except Exception:
                    pass
            except Exception as e:
                logger.error(f"ocr_fallback_vision_failed: {e}")

        return {"pages_ocrd": 1, "merged_text": extracted_text_from_ocr, "vision_fallback": vision_used}

    except Exception as exc:
        logger.error("ocr_all_pages_failed", document_id=document_id, error=str(exc))
        return {"pages_ocrd": 0, "error": str(exc)}


@celery_app.task(
    name="documents.chunk_document",
    bind=True,
    max_retries=2,
    acks_late=True,
    time_limit=1200,
    soft_time_limit=1080,
)
def chunk_document(self, document_id: str, tenant_id: str, strategy: str = "hybrid") -> dict:
    """
    Dividir texto en chunks. Lee el texto desde la BD (campo extracted_text
    escrito por el stage extract_text).
    """
    logger.info("chunk_document_started", document_id=document_id, strategy=strategy)

    session, doc = _get_document_session(document_id)
    if not doc:
        raise ValueError(f"Documento {document_id} no encontrado")

    from ai_platform.services.document_chunker import DocumentChunker

    chunker = DocumentChunker()

    # Leer texto extraido desde la BD
    extracted_text = doc.extracted_text or ""
    if not extracted_text:
        logger.warning(
            "chunk_document_no_text",
            document_id=document_id,
            status=doc.status,
        )
        return {"chunk_count": 0, "strategy": strategy, "warning": "no_text_extracted"}

    if strategy == "hybrid":
        chunks = chunker.chunk_hybrid(
            extracted_text, metadata={"document_id": document_id, "tenant_id": tenant_id}
        )
    elif strategy == "semantic":
        chunks = chunker.chunk_semantic(
            extracted_text, metadata={"document_id": document_id, "tenant_id": tenant_id}
        )
    elif strategy == "fixed":
        chunks = chunker.chunk_fixed(
            extracted_text, metadata={"document_id": document_id, "tenant_id": tenant_id}
        )
    elif strategy == "page":
        chunks = chunker.chunk_page(
            extracted_text, metadata={"document_id": document_id, "tenant_id": tenant_id}
        )
    else:
        chunks = chunker.chunk_hybrid(
            extracted_text, metadata={"document_id": document_id, "tenant_id": tenant_id}
        )

    for chunk in chunks:
        doc_chunk = DocumentChunk(
            tenant_id=tenant_id,
            document_id=document_id,
            chunk_index=chunk.chunk_index,
            level=chunk.level,
            chunk_type=chunk.chunk_type,
            content=chunk.content,
            metadata_json=chunk.metadata,
        )
        session.add(doc_chunk)

    session.commit()

    logger.info(
        "chunk_document_completed",
        document_id=document_id,
        chunk_count=len(chunks),
        strategy=strategy,
        text_length=len(extracted_text),
    )

    return {
        "chunk_count": len(chunks),
        "strategy": strategy,
    }


@celery_app.task(
    name="documents.build_fts_index",
    bind=True,
    max_retries=2,
    acks_late=True,
    time_limit=600,
    soft_time_limit=480,
)
def build_fts_index(self, document_id: str, tenant_id: str) -> dict:
    """
    Construir indice tsvector para busqueda full-text.
    Usa PostgreSQL to_tsvector con configuracion espanola.
    Se basa en los chunks creados en chunk_document stage.
    """
    logger.info("build_fts_index_started", document_id=document_id)

    try:
        session = session_factory()
        try:
            stmt = select(DocumentChunk).where(
                DocumentChunk.document_id == document_id,
                DocumentChunk.tenant_id == tenant_id,
                DocumentChunk.chunk_type == "text",
            )
            chunks = session.execute(stmt).scalars().all()

            fts_entries = []
            for chunk in chunks:
                fts_entry = DocumentFTSIndex(
                    tenant_id=tenant_id,
                    document_id=document_id,
                    chunk_id=chunk.id,
                    chunk_index=chunk.chunk_index,
                    level=chunk.level,
                    search_vector="",
                )
                fts_entries.append(fts_entry)

            # Generar tsvector real usando PostgreSQL nativo
            if fts_entries:
                chunk_ids = [c.id for c in fts_entries]
                raw_sql = """
                    UPDATE document_fts_index
                    SET search_vector = to_tsvector('spanish', content)
                    FROM document_chunks
                    WHERE document_fts_index.chunk_id = document_chunks.id
                      AND document_chunks.id = ANY(:ids)
                """
                session.execute(sa_text(raw_sql), {"ids": chunk_ids})
                session.commit()

            logger.info(
                "build_fts_index_completed",
                document_id=document_id,
                indexed_chunks=len(fts_entries),
            )

            return {"indexed_chunks": len(fts_entries)}

        finally:
            session.close()

    except Exception as exc:
        logger.warning("build_fts_index_failed (optional)", document_id=document_id, error=str(exc))
        return {"indexed_chunks": 0, "warning": str(exc)}


@celery_app.task(
    name="documents.generate_summaries",
    bind=True,
    max_retries=3,
    acks_late=True,
    time_limit=3600,
    soft_time_limit=3300,
)
def generate_summaries(self, document_id: str, tenant_id: str) -> dict:
    """
    Generar resúmenes jerárquicos usando el LLM cliente:
    1. Resúmenes por sección (nivel 2)
    2. Resumen del documento completo (nivel 3)
    """
    logger.info("generate_summaries_started", document_id=document_id)

    session = session_factory()
    try:
        stmt = select(DocumentChunk).where(
            DocumentChunk.document_id == document_id,
            DocumentChunk.tenant_id == tenant_id,
            DocumentChunk.level == 1,
        )
        level1_chunks = session.execute(stmt).scalars().all()

        if not level1_chunks:
            return {"status": "skipped", "reason": "no level 1 chunks found"}

        # Importar LLM client (sincrono, compatible con Celery)
        from ai_platform.orchestrator.llm_client import LLMClient

        llm = LLMClient()
        summaries = {"level2_count": 0, "level3_summaries": []}

        BATCH_SIZE = 5
        level2_results = []

        for i in range(0, len(level1_chunks), BATCH_SIZE):
            batch_chunks = level1_chunks[i : i + BATCH_SIZE]
            batch_text = "\n\n".join(c.content[:2000] for c in batch_chunks)

            prompt = (
                "Eres un asistente que genera resúmenes concisos de texto. "
                f"Resume el siguiente texto en máximo 3 párrafos, conservando los datos clave:\n\n{batch_text}"
            )

            try:
                result = llm.chat(prompt, tenant_id=tenant_id)
                content = result.get("content", "")
                if content and content.strip():
                    level2_summary = DocumentChunk(
                        tenant_id=tenant_id,
                        document_id=document_id,
                        chunk_index=i // BATCH_SIZE,
                        level=2,
                        chunk_type="summary",
                        content=content.strip(),
                        metadata_json={
                            "source_chunk_indices": [c.chunk_index for c in batch_chunks],
                            "source_tenant_id": tenant_id,
                        },
                    )
                    session.add(level2_summary)
                    session.flush()
                    level2_results.append(level2_summary)
                    summaries["level2_count"] += 1
            except Exception as e:
                logger.warning(f"generate_summary_section_failed, batch={i}: {e}")
                continue

        # Generar resumen global del documento (nivel 3)
        if level2_results:
            global_text = "\n\n".join(
                c.content[:2000] for c in level2_results[:10]
            )
            prompt_global = (
                "Eres un asistente que genera un resumen ejecutivo de un documento. "
                "Con base en los siguientes resúmenes de secciones, genera un resumen global "
                "de máximo 5 párrafos que capture los puntos clave del documento:\n\n" + global_text
            )

            try:
                result = llm.chat(prompt_global, tenant_id=tenant_id)
                content = result.get("content", "")
                if content and content.strip():
                    global_summary = DocumentChunk(
                        tenant_id=tenant_id,
                        document_id=document_id,
                        chunk_index=999,
                        level=3,
                        chunk_type="document_summary",
                        content=content.strip(),
                        metadata_json={
                            "source_level2_count": len(level2_results),
                            "source_tenant_id": tenant_id,
                        },
                    )
                    session.add(global_summary)
                    summaries["level3_summaries"].append(content.strip()[:200])
            except Exception as e:
                logger.warning(f"generate_summary_global_failed: {e}")

        summaries["status"] = "completed"
        summaries["total_level2"] = summaries["level2_count"]
        summaries["total_level3"] = len(summaries["level3_summaries"])

        session.commit()

        logger.info(
            "generate_summaries_completed",
            document_id=document_id,
            level2_count=summaries["level2_count"],
            level3_count=len(summaries["level3_summaries"]),
        )

        return summaries

    except Exception as exc:
        logger.error("generate_summaries_failed", document_id=document_id, error=str(exc))
        return {"status": "error", "error": str(exc)}
    finally:
        session.close()
        try:
            llm.close()
        except Exception:
            pass


@celery_app.task(
    name="documents.mark_completed",
    bind=True,
    max_retries=1,
    acks_late=True,
    time_limit=60,
)
def mark_completed(self, document_id: str, tenant_id: str, stats: dict) -> dict:
    """Marcar documento como completado y registrar stats."""
    logger.info("mark_completed", document_id=document_id, stats=stats)

    try:
        _update_document(
            document_id,
            {
                "status": "completed",
                "completed_at": None,
                **stats,
            },
        )
        return {"status": "completed", "document_id": document_id}

    except Exception as exc:
        logger.error("mark_completed_failed", document_id=document_id, error=str(exc))
        raise self.retry(exc=exc, countdown=30) from None


@celery_app.task(
    name="documents.ingest_pipeline",
    bind=True,
    max_retries=1,
    acks_late=True,
    time_limit=7200,
    soft_time_limit=6900,
)
def ingest_pipeline(
    self,
    document_id: str,
    tenant_id: str,
    strategy: str = "hybrid",
) -> dict:
    """
    Pipeline principal de ingestion de documentos.
    """
    logger.info("pipeline_started", document_id=document_id, tenant_id=tenant_id, strategy=strategy)

    try:
        # Stage 1: Extract text
        _update_document(document_id, {"status": "extracting"})
        extract_result = extract_text.delay(document_id, tenant_id)
        extract_data = extract_result.get(timeout=1800)

        # Stage 2: Extract embedded images (DOCX, PDF, PPTX, XLSX)
        _update_document(document_id, {"status": "extracting_images"})
        images_result = extract_embedded_images.delay(document_id, tenant_id)
        images_data = images_result.get(timeout=600)
        images_count = images_data.get("images_extracted", 0)

        # Stage 3: Vision analyze embedded images (async)
        if images_count > 0:
            logger.info("pipeline_vision_needed", document_id=document_id, images=images_count)
            _update_document(document_id, {"status": "vision_analyze"})
            vision_result = vision_analyze_images.delay(document_id, tenant_id)
            _ = vision_result.get(timeout=3600)

        # Stage 4: OCR (conditional)
        if extract_data.get("needs_ocr"):
            logger.info("pipeline_ocr_needed", document_id=document_id)
            _update_document(document_id, {"status": "ocr"})
            ocr_result = ocr_all_pages.delay(document_id, tenant_id, extract_data.get("page_count", 1))
            _ = ocr_result.get(timeout=600)

        # Stage 5: Chunk (lee de la BD, no del Celery result)
        _update_document(document_id, {"status": "chunking"})
        chunk_result = chunk_document.delay(document_id, tenant_id, strategy)
        chunk_data = chunk_result.get(timeout=1200)

        # Stage 6: FTS
        _update_document(document_id, {"status": "indexing"})
        fts_result = build_fts_index.delay(document_id, tenant_id)
        fts_data = fts_result.get(timeout=600)

        # Stage 7: Summaries
        _update_document(document_id, {"status": "summarizing"})
        summary_result = generate_summaries.delay(document_id, tenant_id)
        _ = summary_result.get(timeout=3600)

        # Stage 8: Complete
        mark_completed.delay(
            document_id,
            tenant_id,
            {
                "stats": {
                    "chunks": chunk_data.get("chunk_count", 0),
                    "fts_indexed": fts_data.get("indexed_chunks", 0),
                    "embedded_images": images_count,
                }
            },
        )

        logger.info("pipeline_completed", document_id=document_id)
        return {
            "status": "completed",
            "document_id": document_id,
            "stats": {
                "chunks": chunk_data.get("chunk_count", 0),
                "embedded_images": images_count,
            },
        }

    except Exception as exc:
        logger.error("pipeline_failed", document_id=document_id, error=str(exc))
        _update_document(document_id, {"status": "failed", "error": str(exc)})
        raise self.retry(exc=exc, countdown=120) from None