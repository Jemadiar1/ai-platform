"""
Generador central de documentos profesionales.

Integra:
- Template engine (Jinja2 + FileSystemLoader)
- Markdown a HTML/DOCX/XLSX/PPTX/PDF/PNG
- Markdown a XLSX (multi-tabla)
- Markdown a PPTX (título, secciones, imágenes)

Ejemplo:
    from ai_platform.modules.ai_documents.generators import Generators
    gen = Generators(tenant_id="abc123")
    result = gen.render_docx(tenant_id, {
        "subject": "Marketing Q3 2026",
        "audience": "C-level executives",
        "theme": {"primary_color": "#1a73e8", "company_name": "Acme Inc"},
        "include_chart": True,
    })
    # → result: {"docx": b"PK...", "file_path": "..."}
"""

import base64
import csv
import io
import logging
import os
import time
from dataclasses import asdict, dataclass, field
from typing import Any
from uuid import uuid4

from ai_platform.core.config import get_settings
from ai_platform.database import make_session
from ai_platform.models.db import GeneratedReport, UsageEvent

logger = logging.getLogger(__name__)
settings = get_settings()

_TEMPLATE_DIR = os.path.join(os.path.dirname(__file__), "templates")

# =========================================================================
# Minimal inline template renderer (no external files needed for MVP)
# =========================================================================

_HTML_TEMPLATE = """\
<!DOCTYPE html>
<html lang="es">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{{ title }}</title>
<style>
@page { margin: 2.5cm; }
body {
  font-family: {{ font }}; font-size: 11pt; line-height: 1.6;
  color: {{ font_color }}; background: {{ bg }};
  margin: 0; padding: 0;
}
.page-header {
  background: {{ primary }}; color: white;
  padding: 24px 32px;
}
.page-header h1 { margin: 0 0 4px 0; font-size: 22pt; }
.page-header .meta { opacity: .85; font-size: 9pt; }
.content { padding: 24px 32px; }
.content h2 { color: {{ primary }}; border-bottom: 2px solid {{ primary }}; padding-bottom: 6px; margin-top: 28px; }
.content h3 { color: {{ secondary }}; margin-top: 20px; }
table { border-collapse: collapse; width: 100%; margin: 12px 0; }
th { background: {{ primary }}; color: white; padding: 8px 12px; text-align: left; }
td { border: 1px solid #ddd; padding: 8px 12px; }
tr:nth-child(even) { background: #fafafa; }
img { max-width: 100%; height: auto; margin: 12px 0; text-align: center; display: block; }
.chart-caption { text-align: center; color: {{ secondary }}; font-size: 9pt; margin: -8px 0 12px 0; }
.footer {
  margin-top: 24px; padding: 12px 32px; border-top: 1px solid #ddd;
  font-size: 8pt; color: {{ secondary }}; text-align: center;
}
</style>
</head>
<body>
<div class="page-header">
  <h1>{{ title }}</h1>
  <div class="meta">Audiencia: {{ audience }} &nbsp;|&nbsp; {{ generated_by }}</div>
</div>
<div class="content">
{% for block in blocks %}
  {% if block.type == "heading" %}
    <h2 {{ "style=font-size: 14pt;" if block.level == 1 }}>
      {{ block.text }}
    </h2>
  {% elif block.type == "subheading" %}
    <h3>{{ block.text }}</h3>
  {% elif block.type == "image" %}
    <img src="{{ block.data }}" alt="{{ block.alt or '' }}">
    {% if block.caption %}<div class="chart-caption">{{ block.caption }}</div>{% endif %}
  {% elif block.type == "table" %}
    <table>
      <thead><tr>{% for h in block.headers %}<th>{{ h }}</th>{% endfor %}</tr></thead>
      <tbody>
        {% for row in block.rows %}
        <tr>{% for cell in row %}<td>{{ cell }}</td>{% endfor %}</tr>
        {% endfor %}
      </tbody>
    </table>
  {% else %}
    <p style="white-space:pre-wrap;">{{ block.text }}</p>
  {% endif %}
{% endfor %}
</div>
<div class="footer">
  {{ company }} &nbsp;|&nbsp; {{ title }} &nbsp;|&nbsp; Generado por {{ generated_by }}
</div>
</body>
</html>
"""

_PPTX_SLIDE_BG = """\
<Style>
  <clr>
    <solidFill><srgbClr val="000000"/></solidFill></clr>
  </clr>
  <font sz="14" b="false" kerning="10"><lan>es</lan></font>
  <punct></punct>
</Style>
"""

# =========================================================================
# Parsed markdown line types
# =========================================================================

@dataclass
class Block:
    type: str  # "heading", "subheading", "paragraph", "image", "table"
    text: str = ""
    level: int = 1
    headers: list[str] = field(default_factory=list)
    rows: list[list[str]] = field(default_factory=list)
    data: str = ""  # base64 encoded image
    alt: str = ""
    caption: str = ""


def parse_markdown(md: str) -> list[Block]:
    """Parsear markdown en bloces estructurados."""
    blocks: list[Block] = []
    lines = md.split("\n")

    current_paragraph: list[str] = []

    def _flush_paragraph() -> None:
        nonlocal current_paragraph
        if current_paragraph:
            blocks.append(Block(type="paragraph", text="\n".join(current_paragraph)))
            current_paragraph = []

    i = 0
    while i < len(lines):
        line = lines[i]

        # Image with caption: ```[alt|caption]b64```
        if line.strip().startswith("```") and line.strip().endswith("```"):
            inner = line.strip()[3:-3]
            if "|" in inner and len(inner) > 6:  # likely base64 data
                alt, caption_or_b64 = inner.split("|", 1)
                blocks.append(Block(
                    type="image",
                    data=caption_or_b64,
                    alt=alt.strip(),
                    caption=caption_or_b64 if len(caption_or_b64) < 200 else "",
                ))
                current_paragraph = []
                i += 1
                continue

        # Heading level 1
        if line.startswith("# ") and not line.startswith("## "):
            _flush_paragraph()
            blocks.append(Block(type="heading", text=line[2:].strip(), level=1))
            i += 1
            continue

        # Heading level 2
        if line.startswith("## "):
            _flush_paragraph()
            blocks.append(Block(type="subheading", text=line[3:].strip()))
            i += 1
            continue

        # Horizontal rule
        if set(line.strip()) <= {"-", "_"}:
            _flush_paragraph()
            i += 1
            continue

        # Table detection (all lines have | delimiters)
        if "|" in line and all(
            "|" in lines[j].strip() or set(lines[j].strip()) == set("")
            for j in range(i + 1, min(i + 6, len(lines)))
            if lines[j].strip() and not lines[j].startswith("#")
        ):
            _flush_paragraph()
            headers: list[str] = [h.strip() for h in lines[i].split("|") if h.strip()]
            rows: list[list[str]] = []
            j = i + 1
            # Skip separator line (|---|---|)
            if j < len(lines) and set(lines[j].strip()) <= {"-", "|", ":", " "}:
                j += 1
            while j < len(lines) and j - i < 10 and lines[j].strip() and "|" in lines[j]:
                row = [c.strip() for c in lines[j].split("|") if c.strip()]
                if row:
                    rows.append(row)
                j += 1
            if headers and rows:
                blocks.append(Block(type="table", headers=headers, rows=rows))
                i = j
            else:
                current_paragraph.append(line)
                i += 1
            continue

        current_paragraph.append(line)
        i += 1

    _flush_paragraph()
    return blocks


# =========================================================================
# Template loader
# =========================================================================

def _get_jinja_env():
    """Create Jinja2 environment with FileSystemLoader and custom filters."""
    from jinja2 import BaseLoader, Environment

    env = Environment(
        loader=BaseLoader(),  # FileSystemLoader not available yet; use inline template
        autoescape=True,
    )

    # Custom filters for professional output
    def currency(value: float, symbol: str = "$") -> str:
        return f"{symbol}{value:,.2f}"

    def pct(value: float) -> str:
        return f"{value * 100:.1f}%"

    def truncate(text: str, length: int = 200) -> str:
        return text[:length] + ("..." if len(text) > length else "")

    env.filters["currency"] = currency
    env.filters["pct"] = pct
    env.filters["truncate"] = truncate
    env.filters["safe"] = lambda x: x  # No-op safe for inlined templates

    return env


def _render_html(blocks: list[Block], ctx: dict) -> str:
    """Render HTML from parsed blocks + context."""
    env = _get_jinja_env()
    template = env.from_string(_HTML_TEMPLATE)
    return template.render(**ctx)


# =========================================================================
# Chart generator (matplotlib)
# =========================================================================

def _generate_chart_b64(
    chart_type: str,
    title: str,
    data: list[dict[str, Any]],
    width: int = 600,
    height: int = 400,
    colors: list[str] | None = None,
) -> str:
    """Generate a chart image as base64-encoded PNG."""
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    fig, ax = plt.subplots(figsize=(width / 100, height / 100))

    labels = [d.get("label", "") for d in data]
    values = [d.get("value", 0) for d in data]

    palette = colors or ["#1a73e8", "#e84393", "#00b894", "#fdcb6e", "#6c5ce7", "#e17055", "#00cec9"]

    if chart_type == "bar":
        ax.bar(labels, values, color=palette[: len(labels)])
    elif chart_type == "line":
        ax.plot(labels, values, marker="o", color=palette[0])
        ax.fill_between(labels, values, alpha=0.1)
    elif chart_type == "pie":
        ax.pie(values, labels=labels, colors=palette[: len(labels)], autopct="%1.1f%%")
        ax.set_title(title)
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
        plt.close(fig)
        buf.seek(0)
        return base64.b64encode(buf.getvalue()).decode("utf-8")
    else:
        ax.bar(labels[:10], values[:10], color=palette[: len(labels[:10])])

    ax.set_title(title)
    ax.tick_params(axis="x", rotation=45)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return base64.b64encode(buf.getvalue()).decode("utf-8")


# =========================================================================
# Generators class
# =========================================================================

class Generators:
    """Central generator for multiple file formats."""

    def __init__(self, tenant_id: str) -> None:
        self.tenant_id = tenant_id

    # -------------------------------------------------------------------------
    # Shared: build output dict from params
    # -------------------------------------------------------------------------

    @staticmethod
    def _parse_params(params: dict) -> dict:
        """Parse common params from handler payload."""
        subject = params.get("subject") or params.get("title") or "Documento generado"
        content = params.get("content") or params.get("markdown") or "Documento de ejemplo"

        theme_data = params.get("theme", {})

        return {
            "subject": subject,
            "content": content,
            "audience": params.get("audience", "General"),
            "company": theme_data.get("company_name", "NeuralCrew Labs"),
            "primary_color": theme_data.get("primary_color", "#1a73e8"),
            "secondary_color": theme_data.get("secondary_color", "#5f6368"),
            "font_family": theme_data.get("font_family", "Arial, sans-serif"),
        }

    # -------------------------------------------------------------------------
    # DOCX generator (Phase 3)
    # -------------------------------------------------------------------------

    def render_docx(self, tenant_id: str, params: dict) -> dict:
        """Generate a professional DOCX from markdown content."""
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor

        start_time = time.time()
        data = self._parse_params(params)

        blocks = parse_markdown(data["content"])

        doc = Document()

        # Apply theme
        style = doc.styles["Normal"]
        font = style.font
        font.name = data["font_family"].split(",")[0].strip() or "Calibri"
        font.size = Pt(11)
        font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # Title page heading
        h = doc.add_heading(data["subject"], level=0)
        h.runs[0].font.color.rgb = RGBColor.from_string(data["primary_color"].lstrip("#"))

        doc.add_paragraph(f"Audiencia: {data['audience']}")
        doc.add_paragraph(f"Generado por {data['company']}")

        # Add horizontal rule
        doc.add_paragraph("_" * 80)

        for block in blocks:
            if block.type == "heading":
                h = doc.add_heading(block.text, level=1)
                h.runs[0].font.color.rgb = RGBColor.from_string(data["primary_color"].lstrip("#"))
            elif block.type == "subheading":
                doc.add_heading(block.text, level=2)
            elif block.type == "image":
                try:
                    img_bytes = base64.b64decode(block.data)
                    doc.add_picture(io.BytesIO(img_bytes), width=Inches(5))
                except Exception:
                    doc.add_paragraph(f"[Imagen: {block.alt}]")
            elif block.type == "table":
                table = doc.add_table(
                    rows=1 + len(block.rows),
                    cols=len(block.headers),
                    style="Table Grid",
                )
                # Header row with theme color
                for col_idx, header in enumerate(block.headers):
                    cell = table.rows[0].cells[col_idx]
                    cell.text = str(header)
                    for p in cell.paragraphs:
                        p.runs[0].bold = True
                        p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    cell.paragraphs[0].style = doc.styles["Normal"]
                    # Color header background
                    from docx.oxml.ns import qn
                    shading_elm = cell._element.get_or_add_tcPr()
                    shading = shading_elm.makeelement(qn("w:shd"), {
                        qn("w:fill"): data["primary_color"].lstrip("#"),
                        qn("w:val"): "clear",
                    })
                    shading_elm.append(shading)
                # Data rows
                for row_idx, row_data in enumerate(block.rows, 1):
                    for col_idx, cell_text in enumerate(row_data):
                        table.rows[row_idx].cells[col_idx].text = str(cell_text)
            elif block.type == "paragraph":
                # Parse bold (*text* or **text**) and italic (_text_ or __text__)
                import re
                parts = re.split(r"(\*\*.*?\*\*|[*].*?[*])", block.text)
                if not block.text.strip():
                    doc.add_paragraph()
                    continue
                for part in parts:
                    if part.startswith("**") and part.endswith("**"):
                        p = doc.add_paragraph()
                        r = p.add_run(part[2:-2])
                        r.bold = True
                    elif part.startswith("*") and part.endswith("*"):
                        p = doc.add_paragraph()
                        r = p.add_run(part[1:-1])
                        r.italic = True
                    else:
                        p = doc.add_paragraph(part)

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)

        docx_bytes = buf.getvalue()
        elapsed_ms = int((time.time() - start_time) * 1000)

        return {
            "docx": docx_bytes,
            "format": "docx",
            "filename": f"{data['subject'][:50]}.docx",
            "file_size_bytes": len(docx_bytes),
            "rendering_ms": elapsed_ms,
        }

    # -------------------------------------------------------------------------
    # XLSX generator (Phase 4)
    # -------------------------------------------------------------------------

    def render_xlsx(self, tenant_id: str, params: dict) -> dict:
        """Generate a professional XLSX from markdown content with multi-tab support."""
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
        from openpyxl.styles.colors import Color
        start_time = time.time()
        data = self._parse_params(params)

        blocks = parse_markdown(data["content"])

        wb = Workbook()

        # Style definitions
        header_font = Font(bold=True, color="FFFFFF", name=data["font_family"].split(",")[0].strip() or "Calibri", size=11)
        header_fill = PatternFill(start_color=data["primary_color"].lstrip("#"), end_color=data["primary_color"].lstrip("#"), fill_type="solid")
        title_font = Font(bold=True, size=14, color=Color(data["primary_color"].lstrip("#")))
        thin_border = Border(
            left=Side(style="thin"), right=Side(style="thin"),
            top=Side(style="thin"), bottom=Side(style="thin"),
        )

        ws = wb.active
        ws.title = "Resumen"

        # Title row
        ws.merge_cells("A1:D1")
        title_cell = ws["A1"]
        title_cell.value = data["subject"]
        title_cell.font = Font(bold=True, size=16, color=data["primary_color"].lstrip("#"))
        ws.row_dimensions[1].height = 30

        # Audience row
        ws["A2"].value = f"Audiencia: {data['audience']}"
        ws["A2"].font = Font(italic=True)

        # Data rows
        row_idx = 3
        for block in blocks:
            if block.type == "table":
                # Header row
                for col_idx, header in enumerate(block.headers, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=header)
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.alignment = Alignment(horizontal="center")
                    cell.border = thin_border
                row_idx += 1

                # Data rows
                for row_data in block.rows:
                    for col_idx, cell_text in enumerate(row_data, 1):
                        cell = ws.cell(row=row_idx, column=col_idx, value=str(cell_text))
                        cell.border = thin_border
                    row_idx += 1

                row_idx += 1  # blank row between tables
            elif block.type == "image" and block.alt:
                ws.cell(row=row_idx, column=1, value=f"[Gráfico: {block.alt}]")
                ws.cell(row=row_idx, column=2, value=f"[Imagen: {block.caption}]")
                row_idx += 1
            elif block.type == "paragraph" and block.text.strip():
                ws.cell(row=row_idx, column=1, value=block.text[:200])
                row_idx += 1

        # Auto-size columns (approximate)
        from openpyxl.utils import get_column_letter

        for col_idx in range(1, ws.max_column + 1):
            max_length = 0
            for row in range(1, row_idx + 1):
                cell = ws.cell(row=row, column=col_idx)
                if cell.value:
                    max_length = max(max_length, len(str(cell.value)))
            ws.column_dimensions[get_column_letter(col_idx)].width = min(max(max_length + 2, 10), 50)

        buf = io.BytesIO()
        wb.save(buf)
        buf.seek(0)

        xlsx_bytes = buf.getvalue()
        elapsed_ms = int((time.time() - start_time) * 1000)

        return {
            "xlsx": xlsx_bytes,
            "format": "xlsx",
            "filename": f"{data['subject'][:50]}.xlsx",
            "file_size_bytes": len(xlsx_bytes),
            "rendering_ms": elapsed_ms,
        }

    # -------------------------------------------------------------------------
    # PPTX generator (Phase 6)
    # -------------------------------------------------------------------------

    def render_pptx(self, tenant_id: str, params: dict) -> dict:
        """Generate a professional PPTX presentation from markdown content."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor as PptxRGBColor

        start_time = time.time()
        data = self._parse_params(params)

        blocks = parse_markdown(data["content"])

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        primary = data["primary_color"].lstrip("#")

        def _set_bg(slide, color: str) -> None:
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = PptxRGBColor.from_string(color)

        def _add_textbox(slide, left: float, top: float, width: float, height: float, text: str,
                         font_size: int = 18, bold: bool = False, color: str = "333333",
                         alignment: int = PP_ALIGN.LEFT) -> None:
            txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = PptxRGBColor.from_string(color)
            p.alignment = alignment

        # Slide 1: Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        _set_bg(slide, primary)
        _add_textbox(slide, 0.8, 1.5, 11.7, 2.5, data["subject"], 36, True, "FFFFFF", PP_ALIGN.CENTER)
        _add_textbox(slide, 0.8, 4.2, 11.7, 1.5, f"Audiencia: {data['audience']}", 16, False, "FFFFFF", PP_ALIGN.CENTER)
        _add_textbox(slide, 0.8, 5.2, 11.7, 0.8, f"Generado por {data['company']}", 12, False, "FFFFFF", PP_ALIGN.CENTER)

        # Content slides (max 10 to avoid huge files from markdown)
        for i, block in enumerate(blocks[:10], 1):
            if i > 10:
                break
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(slide, "FFFFFF")

            # Accent bar at top
            accent_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.15))
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = PptxRGBColor.from_string(primary)
            accent_bar.line.fill.background()

            if block.type == "heading":
                _add_textbox(slide, 0.6, 0.5, 12.1, 1, block.text, 28, True, primary)
            elif block.type == "subheading":
                _add_textbox(slide, 0.6, 0.5, 12.1, 0.8, block.text, 22, True, data["secondary_color"].lstrip("#"))
            elif block.type == "paragraph":
                _add_textbox(slide, 0.6, 1.2, 12.1, 4.5, block.text, 14, False, "333333")
            elif block.type == "image":
                try:
                    img_bytes = base64.b64decode(block.data)
                    slide.shapes.add_picture(
                        io.BytesIO(img_bytes),
                        left=Inches(0.6), top=Inches(1.2),
                        width=Inches(12), height=Inches(4.5),
                    )
                    if block.alt:
                        _add_textbox(slide, 0.6, 5.8, 12.1, 0.5, block.alt, 10, False, "666666", PP_ALIGN.CENTER)
                except Exception:
                    _add_textbox(slide, 0.6, 1.2, 12.1, 0.5, f"[Imagen: {block.alt}]", 12, False, "999999")
            elif block.type == "table":
                _add_textbox(slide, 0.6, 0.5, 12.1, 0.6, block.headers[0] if block.headers else "Tabla", 18, True, primary)

                rows = 1 + len(block.rows)
                cols = len(block.headers)
                if rows <= 1 or cols <= 1:
                    rows = max(rows, 2)
                    cols = max(cols, 2)
                left = min(Inches(0.6), Inches(0.2))
                width = min(Inches(12.1), Inches(12.1))
                height = Inches(min(rows * 0.25, 4.5))
                table_shape = slide.shapes.add_table(rows, cols, left, Inches(1.2), width, height)
                table = table_shape.table

                for col_idx, header in enumerate(block.headers):
                    cell = table.cell(0, col_idx)
                    cell.text = header
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(10)
                            run.font.bold = True
                            run.font.color.rgb = PptxRGBColor(0xFF, 0xFF, 0xFF)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = PptxRGBColor.from_string(primary)

                for row_idx, row_data in enumerate(block.rows):
                    for col_idx, cell_text in enumerate(row_data):
                        cell = table.cell(row_idx + 1, col_idx)
                        cell.text = str(cell_text)
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(10)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        pptx_bytes = buf.getvalue()
        elapsed_ms = int((time.time() - start_time) * 1000)

        return {
            "pptx": pptx_bytes,
            "format": "pptx",
            "filename": f"{data['subject'][:50]}.pptx",
            "file_size_bytes": len(pptx_bytes),
            "rendering_ms": elapsed_ms,
        }

    # -------------------------------------------------------------------------
    # PNG generator (Phase 7)
    # -------------------------------------------------------------------------

    def render_png(self, tenant_id: str, params: dict) -> dict:
        """Generate an image (chart, banner, or infographic) from params."""
        from PIL import Image, ImageDraw, ImageFont

        start_time = time.time()
        data = self._parse_params(params)

        chart_type = params.get("chart_type", params.get("type", "bar"))
        charts_data = params.get("charts", [])
        title = params.get("chart_title", data["subject"])

        # Default dimensions for PNG output
        img_width = params.get("width", 800)
        img_height = params.get("height", 600)

        primary = "#" + data["primary_color"].lstrip("#")

        # Try to render as chart if chart data provided
        if charts_data and any(c.get("data") for c in charts_data):
            for chart_info in charts_data:
                data_type = chart_info.get("type", chart_type)
                chart_title = chart_info.get("title", title)
                chart_data = chart_info.get("data", [])
                colors = chart_info.get("colors")

                if chart_data:
                    b64 = _generate_chart_b64(data_type, chart_title, chart_data, img_width, img_height, colors)
                    elapsed_ms = int((time.time() - start_time) * 1000)
                    img_bytes = base64.b64decode(b64)
                    return {
                        "png": img_bytes,
                        "data_uri": f"data:image/png;base64,{b64}",
                        "format": "png",
                        "filename": f"{chart_title[:50]}.png",
                        "file_size_bytes": len(img_bytes),
                        "rendering_ms": elapsed_ms,
                    }

        # Fallback: generate a text banner
        img = Image.new("RGB", (img_width, img_height), color="white")
        draw = ImageDraw.Draw(img)

        try:
            font_title = ImageFont.truetype("arialbd.ttf", 36)
            font_body = ImageFont.truetype("arial.ttf", 24)
        except (OSError, ImportError):
            font_title = ImageFont.load_default()
            font_body = font_title

        # Top accent bar
        draw.rectangle([0, 0, img_width, 12], fill=primary)

        # Title
        bbox = draw.textbbox((40, 80), title, font=font_title)
        tw = bbox[2] - bbox[0]
        draw.text((40 + (img_width - 40 - tw) // 2, 80), title, fill=primary, font=font_title)

        # Subtitle
        subtitle = f"Audiencia: {data['audience']}"
        bbox = draw.textbbox((20, 140), subtitle, font=font_body)
        sw = bbox[2] - bbox[0]
        draw.text((20 + (img_width - 20 - sw) // 2, 140), subtitle, fill="#555555", font=font_body)

        # Content areas
        blocks = parse_markdown(data["content"])[:3]
        y_pos = 220
        for block in blocks:
            if block.type == "paragraph" and block.text.strip():
                draw.text((40, y_pos), block.text[:120], fill="#333333", font=font_body)
                y_pos += 50

        # Footer
        footer_text = f"Generado por {data['company']}"
        bbox = draw.textbbox((20, img_height - 50), footer_text, font=font_body)
        draw.text((20, img_height - 50), footer_text, fill="#999999", font=font_body)
        # Bottom accent bar
        draw.rectangle([0, img_height - 8, img_width, img_height], fill=primary)

        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)

        png_bytes = buf.getvalue()
        elapsed_ms = int((time.time() - start_time) * 1000)

        return {
            "png": png_bytes,
            "format": "png",
            "filename": f"{data['subject'][:50]}.png",
            "file_size_bytes": len(png_bytes),
            "rendering_ms": elapsed_ms,
        }

    # -------------------------------------------------------------------------
    # PDF generator (Phase 5) — reuses existing report_renderer._generate_pdf
    # -------------------------------------------------------------------------

    def render_pdf(self, tenant_id: str, params: dict) -> dict:
        """Generate PDF from parsed markdown blocks using WeasyPrint."""
        start_time = time.time()
        data = self._parse_params(params)
        blocks = parse_markdown(data["content"])

        ctx = {
            "title": data["subject"],
            "audience": data["audience"],
            "company": data["company"],
            "blocks": [asdict(b) for b in blocks],
            "font": data["font_family"].split(",")[0].strip() or "Arial",
            "primary": data["primary_color"].lstrip("#"),
            "secondary": data["secondary_color"].lstrip("#"),
            "font_color": "333333",
            "bg": "ffffff",
            "generated_by": "ai-documents",
        }

        html_content = _render_html(blocks, ctx)

        try:
            from weasyprint import HTML
            pdf_bytes = HTML(string=html_content).write_pdf()
            elapsed_ms = int((time.time() - start_time) * 1000)
            return {
                "pdf": pdf_bytes,
                "format": "pdf",
                "filename": f"{data['subject'][:50]}.pdf",
                "file_size_bytes": len(pdf_bytes),
                "rendering_ms": elapsed_ms,
            }
        except ImportError:
            logger.warning("weasyprint not installed, falling back to existing renderer")
            from ai_platform.services.report_renderer import ReportRendererService
            from ai_platform.services.report_models import ReportFormat, ReportSpec, Section

            sections = [
                Section(id=str(i), title=b.text if b.type in ("heading", "subheading") else f"Section {i}", content=b.text)
                for i, b in enumerate(blocks)
            ]
            from ai_platform.services.report_models import BrandTheme
            spec = ReportSpec(
                title=data["subject"],
                audience=data["audience"],
                sections=sections,
                theme=BrandTheme(
                    primary_color=data["primary_color"],
                    company_name=data["company"],
                ),
            )
            renderer = ReportRendererService()
            outputs = renderer.render(tenant_id, spec, [ReportFormat.PDF])
            pdf_bytes = outputs.get("pdf", b"")
            elapsed_ms = int((time.time() - start_time) * 1000)
            return {
                "pdf": pdf_bytes,
                "format": "pdf",
                "filename": f"{data['subject'][:50]}.pdf",
                "file_size_bytes": len(pdf_bytes),
                "rendering_ms": elapsed_ms,
            }

    # -------------------------------------------------------------------------
    # Render all formats (composite action)
    # -------------------------------------------------------------------------

    def render_all(self, tenant_id: str, params: dict) -> dict:
        """Generate all supported formats in a single call."""
        docx_result = self.render_docx(tenant_id, params)
        xlsx_result = self.render_xlsx(tenant_id, params)
        pptx_result = self.render_pptx(tenant_id, params)

        # PDF may fail if weasyprint missing
        pdf_result = {}
        try:
            pdf_result = self.render_pdf(tenant_id, params)
        except Exception:
            logger.info("PDF generation skipped (weasyprint not available)")

        # PNG may fail if Pillow is missing
        png_result = {}
        try:
            png_result = self.render_png(tenant_id, params)
        except Exception:
            logger.info("PNG generation skipped (Pillow not available)")

        # Flatten all format results so bytes are directly accessible
        formats: dict[str, dict] = {}
        all_formats = {"docx": docx_result, "xlsx": xlsx_result, "pptx": pptx_result, "pdf": pdf_result, "png": png_result}

        for fmt, result in all_formats.items():
            if result.get(fmt):
                formats[fmt] = {
                    "file": result["filename"],
                    "size_bytes": result["file_size_bytes"],
                    "rendering_ms": result.get("rendering_ms", 0),
                }

        # Save to database
        self._persist(tenant_id, params, formats)

        # Return ALL formats with bytes so _send_document_result() can send them
        return {
            "status": "success",
            **formats,  # docx, xlsx, pptx, pdf, png dicts are now at top level with bytes
        }

    # -------------------------------------------------------------------------
    # Persistence (Phase 8)
    # -------------------------------------------------------------------------

    def _persist(self, tenant_id: str, params: dict, formats: dict) -> None:
        """Save generated document metadata to GeneratedReport."""
        try:
            total_size = sum(v["size_bytes"] for v in formats.values())
            rendered_formats = list(formats.keys())
            report_spec = {
                "subject": params.get("subject", params.get("title", "Documento")),
                "content_preview": params.get("content", "")[:500],
                "formats_available": rendered_formats,
                "formats_details": formats,
            }

            with make_session() as db:
                report = GeneratedReport(
                    tenant_id=tenant_id,
                    title=report_spec["subject"],
                    audience=params.get("audience", "General"),
                    generated_formats=rendered_formats,
                    report_spec=report_spec,
                    file_size_bytes=total_size,
                    rendering_time_ms=0,
                )
                db.add(report)
                db.commit()
                logger.info(f"ai-documents persisted: {report.id} for tenant {tenant_id}")
        except Exception as e:
            logger.error(f"Failed to persist ai-documents result: {e}")

    # -------------------------------------------------------------------------
    # Usage logging
    # -------------------------------------------------------------------------

    def _log_usage(self, tenant_id: str, formats: list[str]) -> None:
        """Log usage event for billing."""
        try:
            with make_session() as db:
                event = UsageEvent(
                    tenant_id=tenant_id,
                    module="ai-documents",
                    event_type="document_generation",
                    extra_data={"formats": formats},
                )
                db.add(event)
                db.commit()
        except Exception as e:
            logger.error(f"Failed to log ai-documents usage: {e}")